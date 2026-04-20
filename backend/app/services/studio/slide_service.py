"""Slide deck generation service.

Supports the staged v2 skill-runtime slide-deck pipeline.
"""

from __future__ import annotations

import asyncio
import logging
import re
import shutil
import subprocess
import uuid
from io import BytesIO
from mimetypes import guess_type
from pathlib import Path

from langfuse import observe, propagate_attributes
from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy.exc import SQLAlchemyError

from app.config import settings
from app.models.studio import SlideDeck
from app.schemas.studio import SlideDeckStatus
from app.services.infra.obs_storage import upload_file_to_obs
from app.services.source.source_service import (
    build_combined_content_from_sources,
    fetch_sources,
)
from app.services.studio.studio_status_service import (
    clear_generation_error,
    mark_generation_as_error,
)
from agent.run_skill_prompt import SkillWorkflowRequest, run_skill_workflow

logger = logging.getLogger(__name__)

SLIDE_DECK_V1 = "v1"
SLIDE_DECK_V2 = "v2"
SKILL_WORKFLOW_NAME = "baoyu-slide-deck"
DEFAULT_WORKFLOW_STYLE = "blueprint"
DEFAULT_WORKFLOW_AUDIENCE = "general"
DEFAULT_SLIDE_COUNT = "8"
SLIDE_PREVIEW_MAX_WIDTH = 1280
SLIDE_THUMB_MAX_WIDTH = 320
SUPPORTED_WORKFLOW_STYLES = {
    "blueprint",
    "chalkboard",
    "corporate",
    "minimal",
    "sketch-notes",
    "watercolor",
    "dark-atmospheric",
    "notion",
    "bold-editorial",
    "editorial-infographic",
    "fantasy-animation",
    "intuition-machine",
    "pixel-art",
    "scientific",
    "vector-illustration",
    "vintage",
}
LEGACY_SLIDE_DURATION_TO_SLIDE_COUNT = {
    "shortest": "2",
    "short": "5",
    "default": DEFAULT_SLIDE_COUNT,
}


def _normalize_workflow_style(slide_style: str | None) -> str:
    """Normalize existing slide style values into a skill style id."""
    if not slide_style:
        return DEFAULT_WORKFLOW_STYLE
    normalized = slide_style.strip().lower()
    if not normalized:
        return DEFAULT_WORKFLOW_STYLE
    if normalized in SUPPORTED_WORKFLOW_STYLES:
        return normalized
    logger.warning(
        "Unknown slide style %r, falling back to %s",
        slide_style,
        DEFAULT_WORKFLOW_STYLE,
    )
    return DEFAULT_WORKFLOW_STYLE


def _normalize_workflow_audience(slide_audience: str | None) -> str:
    """Resolve workflow audience from explicit input or default."""
    if slide_audience and slide_audience.strip():
        return slide_audience.strip()
    return DEFAULT_WORKFLOW_AUDIENCE


def _duration_to_slide_count(slide_duration: str | None) -> str:
    """Map the legacy duration choice into a slide count for v2."""
    if not slide_duration:
        return DEFAULT_SLIDE_COUNT
    return LEGACY_SLIDE_DURATION_TO_SLIDE_COUNT.get(
        slide_duration, DEFAULT_SLIDE_COUNT
    )


def _upload_slide_deck_asset(
    file_bytes: bytes,
    *,
    filename: str,
    content_type: str,
) -> str:
    """Upload one slide-deck asset to object storage and return its object key."""
    unique_id = uuid.uuid4().hex[:12]
    object_filename = f"slides/{unique_id}_{filename}"
    return upload_file_to_obs(
        file_content=file_bytes,
        filename=object_filename,
        content_type=content_type,
        cache_control="public, max-age=31536000, immutable",
    )


def _upload_slide_deck_pdf(pdf_bytes: bytes) -> str:
    """Upload PDF bytes to OSS and return object key."""
    return _upload_slide_deck_asset(
        pdf_bytes,
        filename="deck.pdf",
        content_type="application/pdf",
    )


def _normalize_asset_name(value: str) -> str:
    """Normalize one filename stem for generated slide assets."""
    normalized = re.sub(r"[^A-Za-z0-9._-]+", "_", value.strip())
    return normalized.strip("_") or "slide"


def _build_v2_workflow_options(
    slide_deck: SlideDeck,
    *,
    images_only: bool = False,
) -> dict[str, str]:
    """Normalize current slide-deck inputs into the supported skill params."""
    options = {
        "style": _normalize_workflow_style(getattr(slide_deck, "slide_style", None)),
        "audience": _normalize_workflow_audience(
            getattr(slide_deck, "slide_audience", None)
        ),
        "lang": getattr(slide_deck, "slide_language", "简体中文") or "简体中文",
        "slides": _duration_to_slide_count(
            getattr(slide_deck, "slide_duration", "default")
        ),
    }
    if images_only:
        options["images-only"] = "true"
    return options


def _build_v2_source_markdown(
    slide_deck: SlideDeck,
    combined_content: str,
    workflow_options: dict[str, str],
    *,
    focus_topic: str | None = None,
) -> str:
    """Build the markdown source passed into the staged workflow."""
    lines = [
        f"# {slide_deck.title}",
        "",
        "## Slide Deck Request",
        f"- version: {SLIDE_DECK_V2}",
        f"- language: {workflow_options.get('lang', '简体中文')}",
        f"- style: {workflow_options.get('style', DEFAULT_WORKFLOW_STYLE)}",
        f"- audience: {workflow_options.get('audience', DEFAULT_WORKFLOW_AUDIENCE)}",
        f"- slides: {workflow_options.get('slides', DEFAULT_SLIDE_COUNT)}",
    ]
    if focus_topic and focus_topic.strip():
        lines.append(f"- focus_topic: {focus_topic.strip()}")
    custom_prompt = (getattr(slide_deck, "slide_custom_prompt", None) or "").strip()
    if custom_prompt:
        lines.extend(["", "## Additional Instructions", custom_prompt])
    lines.extend(
        [
            "",
            "## Source Content",
            combined_content or "No source content available.",
        ]
    )
    return "\n".join(lines).strip() + "\n"


def _studio_workflow_dir(slide_deck_id: str) -> Path:
    """Return the local workspace directory for one slide-deck workflow run."""
    backend_root = Path(__file__).resolve().parents[3]
    return backend_root / "agent" / "slide_deck" / "studio" / slide_deck_id


def _read_optional_text(path: Path) -> str | None:
    """Read one text artifact if it exists."""
    if not path.exists():
        return None
    return path.read_text(encoding="utf-8")


def _default_suggested_filename(title: str) -> str:
    """Build a stable filename for generated slide deck PDFs."""
    cleaned = re.sub(r"[^\w\-]+", "_", (title or "slide-deck").strip())
    cleaned = cleaned.strip("_") or "slide-deck"
    return f"{cleaned[:80]}.pdf"


def _parse_outline_slides(outline_text: str | None) -> list[dict]:
    """Extract slide titles and filenames from the generated outline."""
    if not outline_text:
        return []

    pattern = re.compile(
        r"^## Slide (?P<number>\d+) of (?P<total>\d+)\n(?P<body>.*?)(?=^## Slide \d+ of \d+|\Z)",
        re.MULTILINE | re.DOTALL,
    )
    slides: list[dict] = []
    for match in pattern.finditer(outline_text):
        body = match.group("body")
        slide_number = int(match.group("number"))
        slide_type = _match_outline_value(body, r"^\*\*Type\*\*:\s*(.+)$")
        filename = _match_outline_value(body, r"^\*\*Filename\*\*:\s*(.+)$")
        title = _match_outline_value(body, r"^Headline:\s*(.+)$")
        if not title:
            title = slide_type or f"Slide {slide_number}"
        slides.append(
            {
                "slide_number": slide_number,
                "title": title.strip(),
                "type": slide_type,
                "filename": filename,
            }
        )
    return slides


def _match_outline_value(body: str, pattern: str) -> str | None:
    """Read the first regex capture from one outline block."""
    match = re.search(pattern, body, re.MULTILINE)
    if match is None:
        return None
    return match.group(1).strip()


def _collect_prompt_artifacts(prompts_dir: Path) -> list[dict[str, str]]:
    """Collect generated prompt files for artifact metadata."""
    if not prompts_dir.exists():
        return []
    artifacts: list[dict[str, str]] = []
    for prompt_file in sorted(prompts_dir.glob("*.md")):
        artifacts.append(
            {
                "name": prompt_file.name,
                "path": str(prompt_file),
            }
        )
    return artifacts


def _collect_image_artifacts(workflow_dir: Path) -> list[dict[str, str]]:
    """Collect generated slide image files for artifact metadata."""
    artifacts: list[dict[str, str]] = []
    for image_file in sorted(workflow_dir.glob("*-slide-*.png")):
        artifacts.append(
            {
                "name": image_file.name,
                "path": str(image_file),
            }
        )
    return artifacts


def _resolve_slide_title(slides: list[dict], index: int) -> str:
    """Return one human-readable slide title."""
    if 0 <= index < len(slides):
        title = (slides[index].get("title") or "").strip()
        if title:
            return title
    return f"Slide {index + 1}"


def _save_image_variant(
    image,
    *,
    output_path: Path,
    max_width: int,
    format_name: str,
    save_kwargs: dict,
) -> tuple[int, int]:
    """Save one resized image variant and return its dimensions."""
    derived = image.copy()
    if derived.width > max_width:
        derived.thumbnail((max_width, max_width * 4))
    derived.save(output_path, format=format_name, **save_kwargs)
    return derived.width, derived.height


def _build_slide_image_variants(
    workflow_dir: Path,
    slides: list[dict],
) -> list[dict]:
    """Create preview/thumb variants and upload all slide image assets."""
    from PIL import Image

    variants_dir = workflow_dir / "variants"
    variants_dir.mkdir(parents=True, exist_ok=True)
    artifacts: list[dict] = []
    for index, export_path in enumerate(sorted(workflow_dir.glob("*-slide-*.png"))):
        with Image.open(export_path) as original:
            width, height = original.size
            safe_stem = _normalize_asset_name(export_path.stem)
            preview_path = variants_dir / f"{safe_stem}-preview.webp"
            thumb_path = variants_dir / f"{safe_stem}-thumb.webp"
            preview_width, preview_height = _save_image_variant(
                original,
                output_path=preview_path,
                max_width=SLIDE_PREVIEW_MAX_WIDTH,
                format_name="WEBP",
                save_kwargs={"quality": 85, "method": 6},
            )
            thumb_width, thumb_height = _save_image_variant(
                original,
                output_path=thumb_path,
                max_width=SLIDE_THUMB_MAX_WIDTH,
                format_name="WEBP",
                save_kwargs={"quality": 78, "method": 6},
            )
        slide_title = _resolve_slide_title(slides, index)
        slide_number = index + 1
        export_key = _upload_slide_deck_asset(
            export_path.read_bytes(),
            filename=f"{safe_stem}-export.png",
            content_type="image/png",
        )
        preview_key = _upload_slide_deck_asset(
            preview_path.read_bytes(),
            filename=f"{safe_stem}-preview.webp",
            content_type="image/webp",
        )
        thumb_key = _upload_slide_deck_asset(
            thumb_path.read_bytes(),
            filename=f"{safe_stem}-thumb.webp",
            content_type="image/webp",
        )
        artifacts.append(
            {
                "index": index,
                "slide_number": slide_number,
                "title": slide_title,
                "filename": export_path.name,
                "variants": {
                    "export": {
                        "filename": export_path.name,
                        "content_type": "image/png",
                        "local_path": str(export_path),
                        "object_key": export_key,
                        "width": width,
                        "height": height,
                    },
                    "preview": {
                        "filename": preview_path.name,
                        "content_type": "image/webp",
                        "local_path": str(preview_path),
                        "object_key": preview_key,
                        "width": preview_width,
                        "height": preview_height,
                    },
                    "thumb": {
                        "filename": thumb_path.name,
                        "content_type": "image/webp",
                        "local_path": str(thumb_path),
                        "object_key": thumb_key,
                        "width": thumb_width,
                        "height": thumb_height,
                    },
                },
            }
        )
    return artifacts


async def _run_bun_merge(
    script_path: Path,
    workflow_dir: Path,
) -> dict[str, str]:
    """Run one slide-deck merge script through bun."""

    def _invoke() -> dict[str, str]:
        commands = [
            ["bun", str(script_path), str(workflow_dir)],
            ["bunx", "bun", str(script_path), str(workflow_dir)],
            ["npx", "-y", "bun", str(script_path), str(workflow_dir)],
        ]
        available_commands = [
            command for command in commands if shutil.which(command[0]) is not None
        ]
        if not available_commands:
            raise RuntimeError(
                "No Bun runtime available for slide merge. "
                "Install `bun` or provide `npx` in the backend runtime."
            )

        command = available_commands[0]
        try:
            completed = subprocess.run(
                command,
                cwd=str(script_path.parent),
                capture_output=True,
                text=True,
                check=True,
            )
        except subprocess.CalledProcessError as exc:
            stdout = (exc.stdout or "").strip()
            stderr = (exc.stderr or "").strip()
            details = [
                f"Slide merge command failed with exit code {exc.returncode}.",
                f"command={' '.join(command)}",
            ]
            if stdout:
                details.append(f"stdout={stdout}")
            if stderr:
                details.append(f"stderr={stderr}")
            raise RuntimeError(" | ".join(details)) from exc
        return {
            "command": " ".join(command),
            "stdout": completed.stdout.strip(),
            "stderr": completed.stderr.strip(),
        }

    return await asyncio.to_thread(_invoke)


async def _merge_slide_deck_outputs(
    workflow_dir: Path,
) -> tuple[Path, Path, dict[str, dict[str, str]]]:
    """Merge generated slide images into PDF and PPTX artifacts."""
    skill_dir = (
        Path(__file__).resolve().parents[3]
        / "agent"
        / "skills"
        / SKILL_WORKFLOW_NAME
    )
    pdf_script = skill_dir / "scripts" / "merge-to-pdf.ts"
    pptx_script = skill_dir / "scripts" / "merge-to-pptx.ts"

    pdf_logs = await _run_bun_merge(pdf_script, workflow_dir)
    pptx_logs = await _run_bun_merge(pptx_script, workflow_dir)

    pdf_path = workflow_dir / f"{workflow_dir.name}.pdf"
    pptx_path = workflow_dir / f"{workflow_dir.name}.pptx"
    if not pdf_path.exists():
        raise FileNotFoundError(f"Merged PDF not found: {pdf_path}")
    if not pptx_path.exists():
        raise FileNotFoundError(f"Merged PPTX not found: {pptx_path}")

    return pdf_path, pptx_path, {"pdf": pdf_logs, "pptx": pptx_logs}


async def _run_slide_deck_generation_v2(
    db: AsyncSession,
    slide_deck: SlideDeck,
    source_ids: list[str] | None,
    *,
    focus_topic: str | None = None,
) -> SlideDeck:
    """Run the staged skill-runtime slide-deck pipeline."""
    sources = await fetch_sources(db, slide_deck.notebook_id, source_ids)
    combined_content = await build_combined_content_from_sources(sources)
    workflow_options = _build_v2_workflow_options(slide_deck)

    backend_root = Path(__file__).resolve().parents[3]
    workflow_dir = _studio_workflow_dir(slide_deck.id)
    workflow_dir.mkdir(parents=True, exist_ok=True)

    source_path = workflow_dir / "source.md"
    source_path.write_text(
        _build_v2_source_markdown(
            slide_deck=slide_deck,
            combined_content=combined_content,
            workflow_options=workflow_options,
            focus_topic=focus_topic,
        ),
        encoding="utf-8",
    )

    results = await run_skill_workflow(
        SkillWorkflowRequest(
            skill_name=SKILL_WORKFLOW_NAME,
            source=str(source_path.relative_to(backend_root)),
            output_dir=str(workflow_dir.relative_to(backend_root)),
            through_stage="image",
            options=workflow_options,
            model=settings.litellm_model,
            temperature=0.3,
            max_completion_tokens=4096,
        ),
        workspace=backend_root,
    )

    pdf_path, pptx_path, merge_logs = await _merge_slide_deck_outputs(workflow_dir)
    object_key = _upload_slide_deck_pdf(pdf_path.read_bytes())

    outline_text = _read_optional_text(workflow_dir / "outline.md")
    slides = _parse_outline_slides(outline_text)
    image_artifacts = await asyncio.to_thread(
        _build_slide_image_variants,
        workflow_dir,
        slides,
    )

    slide_deck.slides_data = {
        "generation_version": SLIDE_DECK_V2,
        "workflow": SKILL_WORKFLOW_NAME,
        "options": workflow_options,
        "slides": slides,
        "artifacts": {
            "analysis": _read_optional_text(workflow_dir / "analysis.md"),
            "outline": outline_text,
            "prompts": _collect_prompt_artifacts(workflow_dir / "prompts"),
            "images": image_artifacts,
            "pdf": str(pdf_path),
            "pptx": str(pptx_path),
            "merge_logs": merge_logs,
        },
        "stage_results": [
            {
                "stage": result.stage,
                "output_path": str(result.output_path),
                "summary": result.llm_result.content.strip(),
            }
            for result in results
        ],
    }
    slide_deck.file_path = object_key
    slide_deck.suggested_filename = _default_suggested_filename(slide_deck.title)
    slide_deck.status = SlideDeckStatus.READY.value
    clear_generation_error(slide_deck)
    await db.flush()
    logger.info(
        "Slide deck %s ready via v2, file_path=%s",
        slide_deck.id,
        object_key,
    )
    return slide_deck

def _new_presentation_without_slides():
    """Return a python-pptx Presentation with all default slides removed."""
    from pptx import Presentation

    prs = Presentation()
    sld_id_lst = prs.slides._sldIdLst
    while len(sld_id_lst) > 0:
        r_id = sld_id_lst[0].rId
        prs.part.drop_rel(r_id)
        del sld_id_lst[0]
    return prs


def _blank_slide_layout(prs):
    """Pick a blank layout from the default template (name or index fallback)."""
    for layout in prs.slide_layouts:
        if layout.name == "Blank":
            return layout
    layouts = prs.slide_layouts
    idx = min(6, len(layouts) - 1)
    return layouts[idx]


def build_slide_pptx_from_pdf_bytes(pdf_bytes: bytes) -> bytes:
    """Rasterize each PDF page and embed as a full-width image in a .pptx."""
    import fitz

    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    try:
        if doc.page_count < 1:
            raise ValueError("PDF has no pages")
        prs = _new_presentation_without_slides()
        blank_layout = _blank_slide_layout(prs)
        for page_index in range(doc.page_count):
            page = doc[page_index]
            mat = fitz.Matrix(2, 2)
            pix = page.get_pixmap(matrix=mat, alpha=False)
            png_stream = BytesIO(pix.tobytes("png"))
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(
                png_stream,
                0,
                0,
                width=prs.slide_width,
            )
        out = BytesIO()
        prs.save(out)
        return out.getvalue()
    finally:
        doc.close()


@observe(name="run_slide_deck_generation_for_existing", as_type="generation")
async def run_slide_deck_generation_for_existing(
    db: AsyncSession,
    slide_deck_id: str,
    source_ids: list[str] | None = None,
    focus_topic: str | None = None,
    user_id: str | None = None,
    session_id: str | None = None,
) -> SlideDeck:
    """Run slide deck generation for an existing pending SlideDeck record."""
    with propagate_attributes(
        user_id=user_id or "",
        session_id=session_id or slide_deck_id,
        metadata={"llm": settings.litellm_model},
    ):
        result = await db.execute(
            select(SlideDeck).where(SlideDeck.id == slide_deck_id)
        )
        slide_deck = result.scalar_one_or_none()
        if slide_deck is None:
            raise ValueError(f"SlideDeck not found: {slide_deck_id}")

        slide_deck.status = SlideDeckStatus.PROCESSING.value
        clear_generation_error(slide_deck)
        await db.flush()

        try:
            return await _run_slide_deck_generation_v2(
                db,
                slide_deck,
                source_ids,
                focus_topic=focus_topic,
            )
        except Exception as exc:
            mark_generation_as_error(
                slide_deck,
                "slide deck generation failed",
                str(exc),
            )
            try:
                await db.flush()
            except SQLAlchemyError:
                logger.exception(
                    "Failed to persist slide deck error state for %s",
                    slide_deck.id,
                )
            raise
