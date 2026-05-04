"""Service for processing source documents: parsing, chunking, embedding."""

import asyncio
import csv
import logging
import time
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path

from fastapi import HTTPException, status
from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession

from app.ai.qwen_asr import transcribe_audio
from app.ai.qwen3_vl_video import understand_video
from app.commons.util import get_image_source_content
from app.config import settings
from app.parsers.bilibili_parser import extract_bilibili_transcript
from app.parsers.web_parser import fetch_web_markdown_via_jina
from app.parsers.youtube_parser import extract_youtube_transcript
from app.models.notebook import Notebook
from app.models.source import Source
from app.services.infra.deep_searcher import call_load_files, call_upload
from app.services.source.source_metadata_skill_service import (
    apply_source_metadata_payload,
    run_source_metadata_skill,
)
from app.services.infra.mineru_client import (
    MinerUClientError,
    apply_asset_urls_to_markdown,
    call_mineru_parse,
    guess_content_type,
)
from app.services.infra.obs_storage import (
    delete_parsed_assets_for_source,
    download_file_from_obs,
    generate_presigned_url,
    get_file_url,
    sources_parsed_prefix,
    upload_bytes_at_key,
    upload_file_to_obs,
)

logger = logging.getLogger(__name__)

# Max characters per source to include in combined content for LLM.
_MAX_CONTENT_PER_SOURCE = 10000

# Upload path stores this for images (see extract_text); not usable for LLM.
_IMAGE_PLACEHOLDER_TEXT = "[Image]"

# Stored until Celery finishes MinerU parsing for PDF uploads.
PDF_SOURCE_PENDING_PLACEHOLDER = "[pdf parsing pending]"

_PROJECT_ROOT = Path(__file__).resolve().parents[4]
_SOURCE_MARKDOWN_ROOT = _PROJECT_ROOT / "files"


def chunk_text(
    text: str, chunk_size: int = 1000, overlap: int = 200
) -> list[dict]:
    """Split text into overlapping chunks with position metadata.

    Returns a list of dicts: {"text": ..., "char_start": ..., "char_end": ...}.
    """
    if not text:
        return []

    chunks = []
    start = 0
    while start < len(text):
        end = min(start + chunk_size, len(text))
        piece = text[start:end].strip()
        if piece:
            chunks.append({
                "text": piece,
                "char_start": start,
                "char_end": end,
            })
        start = end - overlap if end < len(text) else end

    return chunks


def extract_text_by_pages(file_bytes: bytes, file_type: str) -> list[dict]:
    """Extract text with page-level metadata.

    Returns a list of dicts:
      {"page_number": int | None, "text": str}
    For non-paged formats the whole content is returned with page_number=None.
    """
    if file_type == "image":
        return [{"page_number": None, "text": _IMAGE_PLACEHOLDER_TEXT}]

    if file_type in ("txt", "markdown", "csv"):
        return [{"page_number": None, "text": file_bytes.decode("utf-8", errors="replace")}]

    if file_type == "pdf":
        try:
            from io import BytesIO
            from pypdf import PdfReader

            reader = PdfReader(BytesIO(file_bytes))
            pages = []
            for i, page in enumerate(reader.pages, 1):
                text = page.extract_text()
                if text and text.strip():
                    pages.append({"page_number": i, "text": text})
            return pages if pages else [{"page_number": None, "text": ""}]
        except Exception as exc:
            logger.warning("PDF text extraction failed: %s", exc)
            return [{"page_number": None, "text": "[Unable to extract PDF content]"}]

    if file_type == "docx":
        try:
            from io import BytesIO
            from docx import Document

            doc = Document(BytesIO(file_bytes))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            full_text = "\n\n".join(paragraphs)
            return [{"page_number": None, "text": full_text}]
        except Exception as exc:
            logger.warning("DOCX text extraction failed: %s", exc)
            return [{"page_number": None, "text": "[Unable to extract DOCX content]"}]

    if file_type == "pptx":
        try:
            from io import BytesIO
            from pptx import Presentation

            presentation = Presentation(BytesIO(file_bytes))
            slides = []
            for i, slide in enumerate(presentation.slides, 1):
                texts = []
                for shape in slide.shapes:
                    text = getattr(shape, "text", "")
                    if text and text.strip():
                        texts.append(text.strip())
                slide_text = "\n".join(texts).strip()
                if slide_text:
                    slides.append({
                        "page_number": i,
                        "text": slide_text,
                    })
            return slides if slides else [{"page_number": None, "text": ""}]
        except Exception as exc:
            logger.warning("PPTX text extraction failed: %s", exc)
            return [{"page_number": None, "text": "[Unable to extract PPTX content]"}]

    if file_type == "audio":
        return [{"page_number": None, "text": ""}]

    return [{"page_number": None, "text": file_bytes.decode("utf-8", errors="replace")}]


# v1 legacy (kept commented for reference per request)
# async def process_source(db: AsyncSession, source_id: str):
#     """Process a source: chunk the content and generate embeddings.
#
#     Chunks preserve page_number (PDF) and paragraph_index for citation.
#     """
#     result = await db.execute(select(Source).where(Source.id == source_id))
#     source = result.scalar_one_or_none()
#     if source is None:
#         return
#
#     source.status = "processing"
#     await db.flush()
#
#     try:
#         try:
#             delete_by_source_id(source.id)
#         except Exception as exc:
#             logger.warning(
#                 "Milvus delete_by_source_id failed for %s (may be first run): %s",
#                 source.id,
#                 exc,
#             )
#         await db.execute(delete(SourceChunk).where(SourceChunk.source_id == source.id))
#         await db.flush()
#
#         content = source.raw_content or ""
#         if not content:
#             source.status = "error"
#             await db.flush()
#             return
#
#         pages = _split_content_to_pages(content, source.type)
#         all_chunks = _build_chunks_from_pages(pages)
#
#         if not all_chunks:
#             source.status = "error"
#             await db.flush()
#             return
#
#         for i, chunk_info in enumerate(all_chunks):
#             chunk = SourceChunk(
#                 source_id=source.id,
#                 content=chunk_info["text"],
#                 chunk_index=i,
#                 metadata_={
#                     "char_start": chunk_info["char_start"],
#                     "char_end": chunk_info["char_end"],
#                     "page_number": chunk_info.get("page_number"),
#                     "paragraph_index": chunk_info.get("paragraph_index"),
#                 },
#             )
#             db.add(chunk)
#         await db.flush()
#         await db.refresh(source)
#
#         chunk_texts = [c["text"] for c in all_chunks]
#         try:
#             embeddings = await embed_chunks(chunk_texts)
#         except Exception:
#             embeddings = [None] * len(chunk_texts)
#
#         chunks_result = await db.execute(
#             select(SourceChunk)
#             .where(SourceChunk.source_id == source.id)
#             .order_by(SourceChunk.chunk_index)
#         )
#         chunks_list = list(chunks_result.scalars().all())
#         chunk_ids_with_vectors = []
#         for chunk in chunks_list:
#             idx = chunk.chunk_index
#             if idx < len(embeddings) and embeddings[idx] is not None:
#                 chunk_ids_with_vectors.append((str(chunk.id), embeddings[idx]))
#         if chunk_ids_with_vectors:
#             chunk_ids = [c[0] for c in chunk_ids_with_vectors]
#             vectors = [c[1] for c in chunk_ids_with_vectors]
#             notebook_id = source.notebook_id
#             try:
#                 insert_vectors(
#                     chunk_ids=chunk_ids,
#                     source_ids=[source.id] * len(chunk_ids),
#                     notebook_ids=[notebook_id] * len(chunk_ids),
#                     vectors=vectors,
#                 )
#             except Exception as exc:
#                 logger.exception("Milvus insert_vectors failed: %s", exc)
#                 raise
#
#         source.status = "ready"
#         await db.flush()
#
#     except Exception:
#         source.status = "error"
#         await db.flush()
#         raise


def _persist_source_markdown_and_upload(source: Source, content: str) -> str:
    """Write source raw content into local markdown and upload to OSS."""
    t0 = time.perf_counter()
    notebook_dir = _SOURCE_MARKDOWN_ROOT / source.notebook_id
    notebook_dir.mkdir(parents=True, exist_ok=True)
    markdown_path = notebook_dir / f"{source.id}.md"
    markdown_path.write_text(content, encoding="utf-8")
    t_after_write = time.perf_counter()

    object_key = upload_file_to_obs(
        file_content=content.encode("utf-8"),
        filename=f"sources/raw_content_md/{source.notebook_id}/{source.id}.md",
        content_type="text/markdown; charset=utf-8",
    )
    oss_url = get_file_url(object_key)
    total_s = time.perf_counter() - t0
    write_s = t_after_write - t0
    upload_s = time.perf_counter() - t_after_write
    logger.info(
        "source markdown uploaded: source_id=%s local=%s oss_url=%s "
        "timings_s write=%.3f oss_put=%.3f total=%.3f",
        source.id,
        markdown_path,
        oss_url,
        write_s,
        upload_s,
        total_s,
    )
    return oss_url


def _build_pdf_markdown_via_mineru(source: Source) -> str:
    """Sync helper (Celery thread): MinerU → OSS parsed assets → markdown."""
    t_pipeline = time.perf_counter()
    fp = (source.file_path or "").strip()
    if not fp:
        raise MinerUClientError("source has no file_path")

    t_del0 = time.perf_counter()
    try:
        delete_parsed_assets_for_source(source.id)
    except RuntimeError:
        logger.warning(
            "Could not clear prior parsed OSS assets for %s", source.id
        )
    delete_s = time.perf_counter() - t_del0

    t_dl0 = time.perf_counter()
    pdf_bytes = download_file_from_obs(fp)
    download_s = time.perf_counter() - t_dl0

    presigned = None
    t_presign0 = time.perf_counter()
    if not settings.mineru_use_multipart:
        presigned = generate_presigned_url(
            fp,
            expiration=int(settings.mineru_oss_presign_seconds),
        )
    presign_s = time.perf_counter() - t_presign0

    title = source.title or "document.pdf"
    if not title.lower().endswith(".pdf"):
        title = f"{title}.pdf"

    t_mineru0 = time.perf_counter()
    result = call_mineru_parse(
        pdf_presigned_url=presigned,
        pdf_bytes=pdf_bytes if settings.mineru_use_multipart else None,
        original_filename=title,
    )
    mineru_http_s = time.perf_counter() - t_mineru0

    prefix = sources_parsed_prefix(source.id)
    upload_jobs: list[tuple[str, str, bytes, str]] = []
    for rel_path, data in result.files:
        norm = rel_path.replace("\\", "/").lstrip("/")
        if not norm or ".." in norm.split("/"):
            continue
        object_key = f"{prefix}{norm}"
        ctype = guess_content_type(norm)
        upload_jobs.append((norm, object_key, data, ctype))

    path_to_url: dict[str, str] = {}

    def _upload_parsed_asset(
        job: tuple[str, str, bytes, str],
    ) -> tuple[str, str]:
        norm, object_key, data, ctype = job
        upload_bytes_at_key(object_key, data, ctype)
        return norm, get_file_url(object_key)

    t_upload0 = time.perf_counter()
    if upload_jobs:
        max_workers = min(8, max(1, len(upload_jobs)))
        with ThreadPoolExecutor(max_workers=max_workers) as pool:
            futures = [
                pool.submit(_upload_parsed_asset, job)
                for job in upload_jobs
            ]
            for fut in as_completed(futures):
                norm, url = fut.result()
                path_to_url[norm] = url
    oss_parallel_s = time.perf_counter() - t_upload0

    total_s = time.perf_counter() - t_pipeline
    asset_bytes = sum(len(j[2]) for j in upload_jobs)
    logger.info(
        "pdf mineru pipeline timings source_id=%s "
        "delete_s=%.3f download_s=%.3f presign_s=%.3f "
        "mineru_http_s=%.3f oss_parallel_upload_s=%.3f "
        "asset_count=%s asset_bytes=%s total_s=%.3f",
        source.id,
        delete_s,
        download_s,
        presign_s,
        mineru_http_s,
        oss_parallel_s,
        len(upload_jobs),
        asset_bytes,
        total_s,
    )

    return apply_asset_urls_to_markdown(result.markdown, path_to_url)


async def process_source_v2(db: AsyncSession, source_id: str) -> str | None:
    """Process source and persist raw_content as markdown in local + OSS."""
    result = await db.execute(select(Source).where(Source.id == source_id))
    source = result.scalar_one_or_none()
    if source is None:
        return None

    source.status = "processing"
    await db.flush()

    try:
        if source.type == "pdf":
            if not (source.file_path or "").strip():
                source.status = "error"
                source.raw_content = "[PDF] Missing file in object storage."
                await db.flush()
                return None
            if not (settings.mineru_base_url or "").strip():
                source.status = "error"
                source.raw_content = (
                    "[PDF] MinerU is not configured (mineru_base_url)."
                )
                await db.flush()
                return None
            try:
                pdf_md = await asyncio.to_thread(
                    _build_pdf_markdown_via_mineru,
                    source,
                )
            except MinerUClientError as exc:
                logger.warning("MinerU failed for %s: %s", source.id, exc)
                source.status = "error"
                source.raw_content = f"[PDF parsing failed] {exc}"
                await db.flush()
                return None
            stripped = (pdf_md or "").strip()
            if not stripped:
                source.status = "error"
                source.raw_content = "[PDF parsing produced empty content]"
                await db.flush()
                return None
            source.raw_content = stripped
            await db.flush()
            content = stripped
        else:
            content = source.raw_content or ""
            if not content.strip():
                source.status = "error"
                await db.flush()
                return None

        notebook_id = source.notebook_id

        async def _deepsearch_branch() -> str | None:
            markdown_oss_url = await asyncio.to_thread(
                _persist_source_markdown_and_upload,
                source,
                content,
            )
            local_path = await asyncio.to_thread(
                call_upload,
                markdown_oss_url,
                notebook_id,
            )
            if not local_path:
                source.status = "error"
                await db.flush()
                return None

            collection_name = "deepsearcher"
            load_response = await asyncio.to_thread(
                call_load_files,
                base_url=settings.deep_searcher_base_url,
                paths=local_path,
                collection_name=collection_name,
                collection_description="collection desc",
                batch_size=8,
            )
            if not load_response.ok:
                logger.error(
                    "deep load-files failed: source_id=%s status=%s body=%s",
                    source.id,
                    load_response.status_code,
                    (load_response.text or "")[:500],
                )
                source.status = "error"
                await db.flush()
                return None
            try:
                load_body = load_response.json()
            except ValueError:
                load_body = (load_response.text or "")[:500]
            logger.info(
                "deep load-files ok: source_id=%s status=%s body=%s",
                source.id,
                load_response.status_code,
                load_body,
            )
            return markdown_oss_url

        t_gather0 = time.perf_counter()
        deep_out, meta_out = await asyncio.gather(
            _deepsearch_branch(),
            run_source_metadata_skill(
                content,
                source.title,
                source.type,
                log_label=str(source.id),
            ),
            return_exceptions=True,
        )
        gather_wall_s = time.perf_counter() - t_gather0
        logger.info(
            "process_source deepsearch+metadata_gather wall_s=%.3f "
            "source_id=%s",
            gather_wall_s,
            source.id,
        )

        if isinstance(deep_out, Exception):
            logger.exception(
                "deepsearch pipeline failed for %s",
                source.id,
            )
            source.status = "error"
            await db.flush()
            return None

        if deep_out is None:
            return None

        if isinstance(meta_out, Exception):
            logger.exception(
                "Source metadata skill failed for %s",
                source.id,
            )
        elif isinstance(meta_out, dict) and meta_out:
            apply_source_metadata_payload(source, meta_out)

        source.status = "ready"
        await db.flush()
        return deep_out

    except Exception:
        source.status = "error"
        await db.flush()
        raise


async def process_source(db: AsyncSession, source_id: str) -> str | None:
    """v2 entrypoint kept under original method name."""
    return await process_source_v2(db, source_id)


async def finalize_url_source(db: AsyncSession, source: Source) -> None:
    """Fetch text from ``original_url`` for web / YouTube / Bilibili, then chunk.

    On success, persists ``raw_content`` and runs :func:`process_source`.
    """
    url = (source.original_url or "").strip()
    if not url:
        source.status = "error"
        source.raw_content = "[未提供 URL]"
        await db.flush()
        return

    try:
        if source.type == "web":
            text = await fetch_web_markdown_via_jina(url)
        elif source.type == "youtube":
            text = await asyncio.to_thread(extract_youtube_transcript, url)
        elif source.type == "bilibili":
            text = await asyncio.to_thread(extract_bilibili_transcript, url)
        else:
            logger.warning(
                "finalize_url_source unsupported type %s", source.type
            )
            source.status = "error"
            source.raw_content = "[不支持的来源类型]"
            await db.flush()
            return

        stripped = (text or "").strip()
        if not stripped:
            source.status = "error"
            source.raw_content = "[未能从链接提取正文]"
            await db.flush()
            return

        if stripped.startswith("Failed to extract transcript:") or stripped.startswith(
            "Failed to extract Bilibili transcript:"
        ):
            source.status = "error"
            source.raw_content = stripped
            await db.flush()
            return

        if source.type == "bilibili" and stripped.startswith("[Bilibili 字幕]"):
            source.status = "error"
            source.raw_content = stripped
            await db.flush()
            return

        source.raw_content = stripped
        await db.flush()
        await process_source(db, source.id)
    except Exception:
        logger.exception("finalize_url_source failed for %s", source.id)
        source.status = "error"
        source.raw_content = (
            "[链接内容获取失败] 请检查 URL 与网络。"
        )
        await db.flush()


async def finalize_uploaded_video(db: AsyncSession, source: Source) -> None:
    """Run VL on video in OSS, save description to raw_content, then chunk/embed."""
    if source.type != "video" or not source.file_path:
        return
    source.status = "processing"
    await db.flush()
    try:
        video_url = generate_presigned_url(
            source.file_path, expiration=3600
        )
        vl_text = await understand_video(video_url)
    except Exception:
        logger.exception(
            "Video understanding failed for source %s", source.id
        )
        source.status = "error"
        source.raw_content = (
            "[视频理解失败] 请检查对象存储外链与 DashScope 配置。"
        )
        await db.flush()
        return

    stripped = (vl_text or "").strip()
    if not stripped:
        source.status = "error"
        source.raw_content = "[无法从视频中提取有效描述]"
        await db.flush()
        return

    source.raw_content = stripped
    await db.flush()
    try:
        await process_source(db, source.id)
    except Exception as exc:
        logger.warning(
            "process_source failed for video source %s: %s",
            source.id,
            exc,
        )
    await db.refresh(source)


async def finalize_uploaded_audio(db: AsyncSession, source: Source) -> None:
    """Run ASR on audio in OSS, save transcript, then chunk/embed."""
    if source.type != "audio" or not source.file_path:
        return
    source.status = "processing"
    await db.flush()
    try:
        audio_url = generate_presigned_url(
            source.file_path, expiration=3600
        )
        transcript = await transcribe_audio(audio_url)
    except Exception:
        logger.exception(
            "Audio transcription failed for source %s", source.id
        )
        source.status = "error"
        source.raw_content = (
            "[音频转写失败] 请检查对象存储外链与 DashScope ASR 配置。"
        )
        await db.flush()
        return

    stripped = (transcript or "").strip()
    if not stripped:
        source.status = "error"
        source.raw_content = "[无法从音频中提取有效文本]"
        await db.flush()
        return

    source.raw_content = stripped
    await db.flush()
    try:
        await process_source(db, source.id)
    except Exception as exc:
        logger.warning(
            "process_source failed for audio source %s: %s",
            source.id,
            exc,
        )
    await db.refresh(source)


async def finalize_uploaded_image(db: AsyncSession, source: Source) -> None:
    """Run VL on image in OSS, save description to raw_content, then chunk/embed."""
    if source.type != "image" or not source.file_path:
        return
    source.status = "processing"
    await db.flush()
    try:
        vl_text = await get_image_source_content(
            source, max_content=_MAX_CONTENT_PER_SOURCE
        )
    except Exception:
        logger.exception(
            "Image understanding failed for source %s", source.id
        )
        source.status = "error"
        source.raw_content = (
            "[图片理解失败] 请检查对象存储外链与 DashScope 配置。"
        )
        await db.flush()
        return

    stripped = (vl_text or "").strip()
    if not stripped or stripped == "[Image description unavailable]":
        source.status = "error"
        source.raw_content = "[无法从图片中提取有效文本描述]"
        await db.flush()
        return

    source.raw_content = stripped
    await db.flush()
    try:
        await process_source(db, source.id)
    except Exception as exc:
        logger.warning(
            "process_source failed for image source %s: %s",
            source.id,
            exc,
        )
    await db.refresh(source)


def _split_content_to_pages(content: str, source_type: str) -> list[dict]:
    """Split raw_content into page-like segments.

    For plain text stored in raw_content, we treat the whole text as one page.
    Page-level splitting for PDF happens during upload via extract_text_by_pages.
    """
    return [{"page_number": None, "text": content}]


def _build_chunks_from_pages(
    pages: list[dict], chunk_size: int = 1000, overlap: int = 200
) -> list[dict]:
    """Build chunks from pages, preserving page_number in each chunk."""
    all_chunks = []
    paragraph_idx = 0
    for page in pages:
        page_chunks = chunk_text(page["text"], chunk_size, overlap)
        for c in page_chunks:
            c["page_number"] = page.get("page_number")
            c["paragraph_index"] = paragraph_idx
            paragraph_idx += 1
            all_chunks.append(c)
    return all_chunks


def extract_text(file_bytes: bytes, file_type: str) -> str:
    """Extract text content from file bytes based on file type.

    Args:
        file_bytes: Raw file content.
        file_type: Source type (txt, markdown, csv, pdf, docx, pptx, image, audio).

    Returns:
        Extracted text content. For images, returns a placeholder.
    """
    if file_type == "image":
        return _IMAGE_PLACEHOLDER_TEXT

    if file_type == "audio":
        return ""

    if file_type in ("txt", "markdown"):
        return file_bytes.decode("utf-8", errors="replace")

    if file_type == "csv":
        try:
            decoded = file_bytes.decode("utf-8-sig", errors="replace")
            reader = csv.reader(decoded.splitlines())
            rows = []
            for row in reader:
                cleaned = [cell.strip() for cell in row]
                if any(cleaned):
                    rows.append(", ".join(cleaned))
            return "\n".join(rows)
        except Exception as exc:
            logger.warning("CSV text extraction failed: %s", exc)
            return "[Unable to extract CSV content]"

    if file_type == "pdf":
        try:
            from io import BytesIO

            from pypdf import PdfReader

            reader = PdfReader(BytesIO(file_bytes))
            pages = []
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    pages.append(text)
            return "\n\n".join(pages)
        except Exception as exc:
            logger.warning("PDF text extraction failed: %s", exc)
            return "[Unable to extract PDF content]"

    if file_type == "docx":
        try:
            from io import BytesIO

            from docx import Document

            doc = Document(BytesIO(file_bytes))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n\n".join(paragraphs)
        except Exception as exc:
            logger.warning("DOCX text extraction failed: %s", exc)
            return "[Unable to extract DOCX content]"

    if file_type == "pptx":
        try:
            from io import BytesIO

            from pptx import Presentation

            presentation = Presentation(BytesIO(file_bytes))
            slides = []
            for i, slide in enumerate(presentation.slides, 1):
                texts = []
                for shape in slide.shapes:
                    text = getattr(shape, "text", "")
                    if text and text.strip():
                        texts.append(text.strip())
                if texts:
                    slides.append(f"[Slide {i}]\n" + "\n".join(texts))
            return "\n\n".join(slides)
        except Exception as exc:
            logger.warning("PPTX text extraction failed: %s", exc)
            return "[Unable to extract PPTX content]"

    # Fallback: try decoding as text
    return file_bytes.decode("utf-8", errors="replace")


async def _get_single_source_content(source: Source) -> str | None:
    """Get source content from database-extracted raw_content only."""
    raw_content = source.raw_content
    if raw_content is None:
        return None

    # For extracted source types, use persisted raw_content directly.
    if source.type in ("pdf", "image", "video", "audio", "csv", "pptx"):
        return raw_content[:_MAX_CONTENT_PER_SOURCE] if raw_content else None

    return raw_content[:_MAX_CONTENT_PER_SOURCE] if raw_content else None


async def build_combined_content_from_sources(
    sources: list[Source],
) -> str:
    """Build a single combined text from multiple sources for mind map LLM."""
    parts = []
    for source in sources:
        content = await _get_single_source_content(source)
        if not content:
            logger.warning(
                "No content available for source '%s', skipping",
                source.title,
            )
            continue
        logger.info(
            "Source '%s' content length: %s characters",
            source.title,
            len(content),
        )
        logger.info(
            "Source '%s' content preview: %s",
            source.title,
            content[:500],
        )
        parts.append(f"[{source.title}]: {content}")
    return "\n\n".join(parts)


async def fetch_sources(
    db: AsyncSession,
    notebook_id: str,
    source_ids: list[str] | None = None,
) -> list[Source]:
    """Load active sources for a notebook, optionally filtered by source_ids."""
    stmt = select(Source).where(
        Source.notebook_id == notebook_id,
        Source.is_active.is_(True),
    )
    if source_ids:
        stmt = stmt.where(Source.id.in_(source_ids))
    result = await db.execute(stmt)
    return list(result.scalars().all())


async def verify_notebook_access(
    db: AsyncSession, notebook_id: str, user_id: str
):
    """Verify the user has access to the notebook."""
    result = await db.execute(
        select(Notebook).where(
            Notebook.id == notebook_id, Notebook.user_id == user_id
        )
    )
    if result.scalar_one_or_none() is None:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Notebook not found",
        )


async def get_source(
    db: AsyncSession, source_id: str, user_id: str
) -> Source:
    """Get a source and verify user access through its notebook."""
    result = await db.execute(
        select(Source)
        .join(Notebook, Source.notebook_id == Notebook.id)
        .where(Source.id == source_id, Notebook.user_id == user_id)
    )
    source = result.scalar_one_or_none()
    if source is None:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Source not found",
        )
    return source


async def get_source_in_notebook(
    db: AsyncSession, source_id: str, notebook_id: str
) -> Source:
    """Get a source that belongs to the given notebook (public share)."""
    result = await db.execute(
        select(Source).where(
            Source.id == source_id,
            Source.notebook_id == notebook_id,
        )
    )
    source = result.scalar_one_or_none()
    if source is None:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Source not found",
        )
    return source
