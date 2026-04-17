"""Service for processing source documents: parsing, chunking, embedding."""

import asyncio
import csv
import logging
from pathlib import Path

from fastapi import HTTPException, status
from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession

from app.ai.qwen_asr import transcribe_audio
from app.ai.qwen3_vl_video import understand_video
from app.commons.util import get_image_source_content
from app.parsers.bilibili_parser import extract_bilibili_transcript
from app.parsers.web_parser import parse_web_page
from app.parsers.youtube_parser import extract_youtube_transcript
from app.models.notebook import Notebook
from app.models.source import Source
from app.services.obs_storage import (
    generate_presigned_url,
    get_file_url,
    upload_file_to_obs,
)

logger = logging.getLogger(__name__)

# Max characters per source to include in combined content for LLM.
_MAX_CONTENT_PER_SOURCE = 10000

# Upload path stores this for images (see extract_text); not usable for LLM.
_IMAGE_PLACEHOLDER_TEXT = "[Image]"

_PROJECT_ROOT = Path(__file__).resolve().parents[3]
_SOURCE_MARKDOWN_ROOT = _PROJECT_ROOT / "files"


def chunk_text(
    text: str, chunk_size: int = 1000, overlap: int = 200
) -> list[dict]:
    """Split text into overlapping chunks with position metadata.

    Returns a list of dicts: {"text": ..., "char_start": ..., "char_end": ...}.
    """
    if not text:
        return []

    chunks = []
    start = 0
    while start < len(text):
        end = min(start + chunk_size, len(text))
        piece = text[start:end].strip()
        if piece:
            chunks.append({
                "text": piece,
                "char_start": start,
                "char_end": end,
            })
        start = end - overlap if end < len(text) else end

    return chunks


def extract_text_by_pages(file_bytes: bytes, file_type: str) -> list[dict]:
    """Extract text with page-level metadata.

    Returns a list of dicts:
      {"page_number": int | None, "text": str}
    For non-paged formats the whole content is returned with page_number=None.
    """
    if file_type == "image":
        return [{"page_number": None, "text": _IMAGE_PLACEHOLDER_TEXT}]

    if file_type in ("txt", "markdown", "csv"):
        return [{"page_number": None, "text": file_bytes.decode("utf-8", errors="replace")}]

    if file_type == "pdf":
        try:
            from io import BytesIO
            from pypdf import PdfReader

            reader = PdfReader(BytesIO(file_bytes))
            pages = []
            for i, page in enumerate(reader.pages, 1):
                text = page.extract_text()
                if text and text.strip():
                    pages.append({"page_number": i, "text": text})
            return pages if pages else [{"page_number": None, "text": ""}]
        except Exception as exc:
            logger.warning("PDF text extraction failed: %s", exc)
            return [{"page_number": None, "text": "[Unable to extract PDF content]"}]

    if file_type == "docx":
        try:
            from io import BytesIO
            from docx import Document

            doc = Document(BytesIO(file_bytes))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            full_text = "\n\n".join(paragraphs)
            return [{"page_number": None, "text": full_text}]
        except Exception as exc:
            logger.warning("DOCX text extraction failed: %s", exc)
            return [{"page_number": None, "text": "[Unable to extract DOCX content]"}]

    if file_type == "pptx":
        try:
            from io import BytesIO
            from pptx import Presentation

            presentation = Presentation(BytesIO(file_bytes))
            slides = []
            for i, slide in enumerate(presentation.slides, 1):
                texts = []
                for shape in slide.shapes:
                    text = getattr(shape, "text", "")
                    if text and text.strip():
                        texts.append(text.strip())
                slide_text = "\n".join(texts).strip()
                if slide_text:
                    slides.append({
                        "page_number": i,
                        "text": slide_text,
                    })
            return slides if slides else [{"page_number": None, "text": ""}]
        except Exception as exc:
            logger.warning("PPTX text extraction failed: %s", exc)
            return [{"page_number": None, "text": "[Unable to extract PPTX content]"}]

    if file_type == "audio":
        return [{"page_number": None, "text": ""}]

    return [{"page_number": None, "text": file_bytes.decode("utf-8", errors="replace")}]


# v1 legacy (kept commented for reference per request)
# async def process_source(db: AsyncSession, source_id: str):
#     """Process a source: chunk the content and generate embeddings.
#
#     Chunks preserve page_number (PDF) and paragraph_index for citation.
#     """
#     result = await db.execute(select(Source).where(Source.id == source_id))
#     source = result.scalar_one_or_none()
#     if source is None:
#         return
#
#     source.status = "processing"
#     await db.flush()
#
#     try:
#         try:
#             delete_by_source_id(source.id)
#         except Exception as exc:
#             logger.warning(
#                 "Milvus delete_by_source_id failed for %s (may be first run): %s",
#                 source.id,
#                 exc,
#             )
#         await db.execute(delete(SourceChunk).where(SourceChunk.source_id == source.id))
#         await db.flush()
#
#         content = source.raw_content or ""
#         if not content:
#             source.status = "error"
#             await db.flush()
#             return
#
#         pages = _split_content_to_pages(content, source.type)
#         all_chunks = _build_chunks_from_pages(pages)
#
#         if not all_chunks:
#             source.status = "error"
#             await db.flush()
#             return
#
#         for i, chunk_info in enumerate(all_chunks):
#             chunk = SourceChunk(
#                 source_id=source.id,
#                 content=chunk_info["text"],
#                 chunk_index=i,
#                 metadata_={
#                     "char_start": chunk_info["char_start"],
#                     "char_end": chunk_info["char_end"],
#                     "page_number": chunk_info.get("page_number"),
#                     "paragraph_index": chunk_info.get("paragraph_index"),
#                 },
#             )
#             db.add(chunk)
#         await db.flush()
#         await db.refresh(source)
#
#         chunk_texts = [c["text"] for c in all_chunks]
#         try:
#             embeddings = await embed_chunks(chunk_texts)
#         except Exception:
#             embeddings = [None] * len(chunk_texts)
#
#         chunks_result = await db.execute(
#             select(SourceChunk)
#             .where(SourceChunk.source_id == source.id)
#             .order_by(SourceChunk.chunk_index)
#         )
#         chunks_list = list(chunks_result.scalars().all())
#         chunk_ids_with_vectors = []
#         for chunk in chunks_list:
#             idx = chunk.chunk_index
#             if idx < len(embeddings) and embeddings[idx] is not None:
#                 chunk_ids_with_vectors.append((str(chunk.id), embeddings[idx]))
#         if chunk_ids_with_vectors:
#             chunk_ids = [c[0] for c in chunk_ids_with_vectors]
#             vectors = [c[1] for c in chunk_ids_with_vectors]
#             notebook_id = source.notebook_id
#             try:
#                 insert_vectors(
#                     chunk_ids=chunk_ids,
#                     source_ids=[source.id] * len(chunk_ids),
#                     notebook_ids=[notebook_id] * len(chunk_ids),
#                     vectors=vectors,
#                 )
#             except Exception as exc:
#                 logger.exception("Milvus insert_vectors failed: %s", exc)
#                 raise
#
#         source.status = "ready"
#         await db.flush()
#
#     except Exception:
#         source.status = "error"
#         await db.flush()
#         raise


def _persist_source_markdown_and_upload(source: Source, content: str) -> str:
    """Write source raw content into local markdown and upload to OSS."""
    notebook_dir = _SOURCE_MARKDOWN_ROOT / source.notebook_id
    notebook_dir.mkdir(parents=True, exist_ok=True)
    markdown_path = notebook_dir / f"{source.id}.md"
    markdown_path.write_text(content, encoding="utf-8")

    object_key = upload_file_to_obs(
        file_content=content.encode("utf-8"),
        filename=f"sources/raw_content_md/{source.notebook_id}/{source.id}.md",
        content_type="text/markdown; charset=utf-8",
    )
    oss_url = get_file_url(object_key)
    logger.info(
        "source markdown uploaded: source_id=%s local=%s oss_url=%s",
        source.id,
        markdown_path,
        oss_url,
    )
    return oss_url


async def process_source_v2(db: AsyncSession, source_id: str) -> str | None:
    """Process source and persist raw_content as markdown in local + OSS."""
    result = await db.execute(select(Source).where(Source.id == source_id))
    source = result.scalar_one_or_none()
    if source is None:
        return None

    source.status = "processing"
    await db.flush()

    try:
        content = source.raw_content or ""
        if not content:
            source.status = "error"
            await db.flush()
            return None

        markdown_oss_url = _persist_source_markdown_and_upload(source, content)

        source.status = "ready"
        await db.flush()
        return markdown_oss_url

    except Exception:
        source.status = "error"
        await db.flush()
        raise


async def process_source(db: AsyncSession, source_id: str) -> str | None:
    """v2 entrypoint kept under original method name."""
    return await process_source_v2(db, source_id)


async def finalize_url_source(db: AsyncSession, source: Source) -> None:
    """Fetch text from ``original_url`` for web / YouTube / Bilibili, then chunk.

    On success, persists ``raw_content`` and runs :func:`process_source`.
    """
    url = (source.original_url or "").strip()
    if not url:
        source.status = "error"
        source.raw_content = "[未提供 URL]"
        await db.flush()
        return

    try:
        if source.type == "web":
            text = await parse_web_page(url)
        elif source.type == "youtube":
            text = await asyncio.to_thread(extract_youtube_transcript, url)
        elif source.type == "bilibili":
            text = await asyncio.to_thread(extract_bilibili_transcript, url)
        else:
            logger.warning(
                "finalize_url_source unsupported type %s", source.type
            )
            source.status = "error"
            source.raw_content = "[不支持的来源类型]"
            await db.flush()
            return

        stripped = (text or "").strip()
        if not stripped:
            source.status = "error"
            source.raw_content = "[未能从链接提取正文]"
            await db.flush()
            return

        if stripped.startswith("Failed to extract transcript:") or stripped.startswith(
            "Failed to extract Bilibili transcript:"
        ):
            source.status = "error"
            source.raw_content = stripped
            await db.flush()
            return

        if source.type == "bilibili" and stripped.startswith("[Bilibili 字幕]"):
            source.status = "error"
            source.raw_content = stripped
            await db.flush()
            return

        source.raw_content = stripped
        await db.flush()
        await process_source(db, source.id)
    except Exception:
        logger.exception("finalize_url_source failed for %s", source.id)
        source.status = "error"
        source.raw_content = (
            "[链接内容获取失败] 请检查 URL 与网络。"
        )
        await db.flush()


async def finalize_uploaded_video(db: AsyncSession, source: Source) -> None:
    """Run VL on video in OSS, save description to raw_content, then chunk/embed."""
    if source.type != "video" or not source.file_path:
        return
    source.status = "processing"
    await db.flush()
    try:
        video_url = generate_presigned_url(
            source.file_path, expiration=3600
        )
        vl_text = await understand_video(video_url)
    except Exception:
        logger.exception(
            "Video understanding failed for source %s", source.id
        )
        source.status = "error"
        source.raw_content = (
            "[视频理解失败] 请检查对象存储外链与 DashScope 配置。"
        )
        await db.flush()
        return

    stripped = (vl_text or "").strip()
    if not stripped:
        source.status = "error"
        source.raw_content = "[无法从视频中提取有效描述]"
        await db.flush()
        return

    source.raw_content = stripped
    await db.flush()
    try:
        await process_source(db, source.id)
    except Exception as exc:
        logger.warning(
            "process_source failed for video source %s: %s",
            source.id,
            exc,
        )
    await db.refresh(source)


async def finalize_uploaded_audio(db: AsyncSession, source: Source) -> None:
    """Run ASR on audio in OSS, save transcript, then chunk/embed."""
    if source.type != "audio" or not source.file_path:
        return
    source.status = "processing"
    await db.flush()
    try:
        audio_url = generate_presigned_url(
            source.file_path, expiration=3600
        )
        transcript = await transcribe_audio(audio_url)
    except Exception:
        logger.exception(
            "Audio transcription failed for source %s", source.id
        )
        source.status = "error"
        source.raw_content = (
            "[音频转写失败] 请检查对象存储外链与 DashScope ASR 配置。"
        )
        await db.flush()
        return

    stripped = (transcript or "").strip()
    if not stripped:
        source.status = "error"
        source.raw_content = "[无法从音频中提取有效文本]"
        await db.flush()
        return

    source.raw_content = stripped
    await db.flush()
    try:
        await process_source(db, source.id)
    except Exception as exc:
        logger.warning(
            "process_source failed for audio source %s: %s",
            source.id,
            exc,
        )
    await db.refresh(source)


async def finalize_uploaded_image(db: AsyncSession, source: Source) -> None:
    """Run VL on image in OSS, save description to raw_content, then chunk/embed."""
    if source.type != "image" or not source.file_path:
        return
    source.status = "processing"
    await db.flush()
    try:
        vl_text = await get_image_source_content(
            source, max_content=_MAX_CONTENT_PER_SOURCE
        )
    except Exception:
        logger.exception(
            "Image understanding failed for source %s", source.id
        )
        source.status = "error"
        source.raw_content = (
            "[图片理解失败] 请检查对象存储外链与 DashScope 配置。"
        )
        await db.flush()
        return

    stripped = (vl_text or "").strip()
    if not stripped or stripped == "[Image description unavailable]":
        source.status = "error"
        source.raw_content = "[无法从图片中提取有效文本描述]"
        await db.flush()
        return

    source.raw_content = stripped
    await db.flush()
    try:
        await process_source(db, source.id)
    except Exception as exc:
        logger.warning(
            "process_source failed for image source %s: %s",
            source.id,
            exc,
        )
    await db.refresh(source)


def _split_content_to_pages(content: str, source_type: str) -> list[dict]:
    """Split raw_content into page-like segments.

    For plain text stored in raw_content, we treat the whole text as one page.
    Page-level splitting for PDF happens during upload via extract_text_by_pages.
    """
    return [{"page_number": None, "text": content}]


def _build_chunks_from_pages(
    pages: list[dict], chunk_size: int = 1000, overlap: int = 200
) -> list[dict]:
    """Build chunks from pages, preserving page_number in each chunk."""
    all_chunks = []
    paragraph_idx = 0
    for page in pages:
        page_chunks = chunk_text(page["text"], chunk_size, overlap)
        for c in page_chunks:
            c["page_number"] = page.get("page_number")
            c["paragraph_index"] = paragraph_idx
            paragraph_idx += 1
            all_chunks.append(c)
    return all_chunks


def extract_text(file_bytes: bytes, file_type: str) -> str:
    """Extract text content from file bytes based on file type.

    Args:
        file_bytes: Raw file content.
        file_type: Source type (txt, markdown, csv, pdf, docx, pptx, image, audio).

    Returns:
        Extracted text content. For images, returns a placeholder.
    """
    if file_type == "image":
        return _IMAGE_PLACEHOLDER_TEXT

    if file_type == "audio":
        return ""

    if file_type in ("txt", "markdown"):
        return file_bytes.decode("utf-8", errors="replace")

    if file_type == "csv":
        try:
            decoded = file_bytes.decode("utf-8-sig", errors="replace")
            reader = csv.reader(decoded.splitlines())
            rows = []
            for row in reader:
                cleaned = [cell.strip() for cell in row]
                if any(cleaned):
                    rows.append(", ".join(cleaned))
            return "\n".join(rows)
        except Exception as exc:
            logger.warning("CSV text extraction failed: %s", exc)
            return "[Unable to extract CSV content]"

    if file_type == "pdf":
        try:
            from io import BytesIO

            from pypdf import PdfReader

            reader = PdfReader(BytesIO(file_bytes))
            pages = []
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    pages.append(text)
            return "\n\n".join(pages)
        except Exception as exc:
            logger.warning("PDF text extraction failed: %s", exc)
            return "[Unable to extract PDF content]"

    if file_type == "docx":
        try:
            from io import BytesIO

            from docx import Document

            doc = Document(BytesIO(file_bytes))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n\n".join(paragraphs)
        except Exception as exc:
            logger.warning("DOCX text extraction failed: %s", exc)
            return "[Unable to extract DOCX content]"

    if file_type == "pptx":
        try:
            from io import BytesIO

            from pptx import Presentation

            presentation = Presentation(BytesIO(file_bytes))
            slides = []
            for i, slide in enumerate(presentation.slides, 1):
                texts = []
                for shape in slide.shapes:
                    text = getattr(shape, "text", "")
                    if text and text.strip():
                        texts.append(text.strip())
                if texts:
                    slides.append(f"[Slide {i}]\n" + "\n".join(texts))
            return "\n\n".join(slides)
        except Exception as exc:
            logger.warning("PPTX text extraction failed: %s", exc)
            return "[Unable to extract PPTX content]"

    # Fallback: try decoding as text
    return file_bytes.decode("utf-8", errors="replace")


async def _get_single_source_content(source: Source) -> str | None:
    """Get source content from database-extracted raw_content only."""
    raw_content = source.raw_content
    if raw_content is None:
        return None

    # For extracted source types, use persisted raw_content directly.
    if source.type in ("pdf", "image", "video", "audio", "csv", "pptx"):
        return raw_content[:_MAX_CONTENT_PER_SOURCE] if raw_content else None

    return raw_content[:_MAX_CONTENT_PER_SOURCE] if raw_content else None


async def build_combined_content_from_sources(
    sources: list[Source],
) -> str:
    """Build a single combined text from multiple sources for mind map LLM."""
    parts = []
    for source in sources:
        content = await _get_single_source_content(source)
        if not content:
            logger.warning(
                "No content available for source '%s', skipping",
                source.title,
            )
            continue
        logger.info(
            "Source '%s' content length: %s characters",
            source.title,
            len(content),
        )
        logger.info(
            "Source '%s' content preview: %s",
            source.title,
            content[:500],
        )
        parts.append(f"[{source.title}]: {content}")
    return "\n\n".join(parts)


async def fetch_sources(
    db: AsyncSession,
    notebook_id: str,
    source_ids: list[str] | None = None,
) -> list[Source]:
    """Load active sources for a notebook, optionally filtered by source_ids."""
    stmt = select(Source).where(
        Source.notebook_id == notebook_id,
        Source.is_active.is_(True),
    )
    if source_ids:
        stmt = stmt.where(Source.id.in_(source_ids))
    result = await db.execute(stmt)
    return list(result.scalars().all())


async def verify_notebook_access(
    db: AsyncSession, notebook_id: str, user_id: str
):
    """Verify the user has access to the notebook."""
    result = await db.execute(
        select(Notebook).where(
            Notebook.id == notebook_id, Notebook.user_id == user_id
        )
    )
    if result.scalar_one_or_none() is None:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Notebook not found",
        )


async def get_source(
    db: AsyncSession, source_id: str, user_id: str
) -> Source:
    """Get a source and verify user access through its notebook."""
    result = await db.execute(
        select(Source)
        .join(Notebook, Source.notebook_id == Notebook.id)
        .where(Source.id == source_id, Notebook.user_id == user_id)
    )
    source = result.scalar_one_or_none()
    if source is None:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Source not found",
        )
    return source


async def get_source_in_notebook(
    db: AsyncSession, source_id: str, notebook_id: str
) -> Source:
    """Get a source that belongs to the given notebook (public share)."""
    result = await db.execute(
        select(Source).where(
            Source.id == source_id,
            Source.notebook_id == notebook_id,
        )
    )
    source = result.scalar_one_or_none()
    if source is None:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Source not found",
        )
    return source
